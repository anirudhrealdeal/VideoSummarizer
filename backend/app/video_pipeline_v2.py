import os
import subprocess
import json
import re
import math
import textwrap
from pathlib import Path
from typing import List, Dict

from PIL import Image, ImageDraw, ImageFont
from openai import OpenAI
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


def _get_client():
    """Get OpenAI client."""
    api_key = os.getenv("OPENAI_API_KEY")
    if not api_key:
        raise EnvironmentError("OPENAI_API_KEY is required in environment")
    return OpenAI(api_key=api_key)


def _run_cmd(cmd: List[str]):
    completed = subprocess.run(cmd, capture_output=True, text=True)
    if completed.returncode != 0:
        raise RuntimeError(f"Command failed: {cmd}\nstderr: {completed.stderr}")
    return completed.stdout


def download_youtube(youtube_url: str, out_dir: str) -> str:
    out_template = os.path.join(out_dir, "video.%(ext)s")
    cmd = [
        "yt-dlp",
        "--no-playlist",
        "-f", "bestvideo[ext=mp4]+bestaudio/best[ext=mp4]/best",
        "-o", out_template,
        youtube_url,
    ]
    _run_cmd(cmd)
    files = list(Path(out_dir).glob("video.*"))
    if not files:
        raise FileNotFoundError("Downloaded video not found")
    return str(files[0])


def extract_audio(video_path: str, out_dir: str) -> str:
    audio_path = os.path.join(out_dir, "audio.wav")
    _run_cmd(["ffmpeg", "-y", "-i", video_path, "-vn", "-ac", "1", "-ar", "16000", audio_path])
    return audio_path


def split_audio(audio_path: str, out_dir: str, chunk_duration: int = 600) -> List[str]:
    """Split audio into chunks of chunk_duration seconds."""
    duration = float(_run_cmd(["ffprobe", "-v", "error", "-show_entries", "format=duration", "-of", "default=noprint_wrappers=1:nokey=1", audio_path]).strip())
    num_chunks = math.ceil(duration / chunk_duration)
    chunk_paths = []
    for i in range(num_chunks):
        start = i * chunk_duration
        chunk_length = min(chunk_duration, duration - start)
        chunk_path = os.path.join(out_dir, f"chunk_{i:03d}.wav")
        _run_cmd(["ffmpeg", "-y", "-i", audio_path, "-ss", str(start), "-t", str(chunk_length), "-c", "copy", chunk_path])
        chunk_paths.append(chunk_path)
    return chunk_paths


def video_duration(video_path: str) -> float:
    out = _run_cmd(["ffprobe", "-v", "error", "-show_entries", "format=duration", "-of", "default=noprint_wrappers=1:nokey=1", video_path])
    return float(out.strip())


def _transcribe_chunk(chunk_path: str) -> str:
    client = _get_client()
    with open(chunk_path, "rb") as f:
        result = client.audio.transcriptions.create(file=f, model="whisper-1")
    return result.text


def speech_to_text(audio_path: str) -> str:
    from concurrent.futures import ThreadPoolExecutor, as_completed
    out_dir = os.path.dirname(audio_path)
    chunk_paths = split_audio(audio_path, out_dir)
    results = [None] * len(chunk_paths)
    with ThreadPoolExecutor(max_workers=6) as executor:
        futures = {executor.submit(_transcribe_chunk, p): i for i, p in enumerate(chunk_paths)}
        for future in as_completed(futures):
            results[futures[future]] = future.result()
    return " ".join(r for r in results if r).strip()


def chunk_transcript_by_time(transcript: str, chunk_words: int = 500) -> List[Dict]:
    """Break transcript into semantic chunks with approximate timestamps."""
    words = transcript.split()
    chunks = []
    words_per_minute = 150  # average speech rate
    
    for i in range(0, len(words), chunk_words):
        chunk_text = " ".join(words[i : i + chunk_words])
        start_time = (i / len(words)) * (len(words) / words_per_minute) * 60 if len(words) > 0 else 0
        chunks.append({
            "text": chunk_text,
            "start_time": int(start_time),
            "index": len(chunks)
        })
    return chunks


def extract_executive_summary(transcript: str) -> str:
    """Generate 1-2 min executive summary."""
    client = _get_client()
    prompt = (
        "You are an expert educational content summarizer. "
        "Generate a VERY CONCISE executive summary (100-150 words, ~1 minute read) "
        "of this video transcript. Focus on the main takeaway and critical insights only.\n\n"
        f"Transcript:\n{transcript[:3000]}...\n\n"
        "Executive Summary:"
    )
    resp = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[{"role": "user", "content": prompt}],
        max_tokens=200
    )
    return resp.choices[0].message.content.strip()


def extract_key_points(transcript: str) -> List[str]:
    """Extract 5-7 key bullet points."""
    client = _get_client()
    prompt = (
        "You are an expert content analyst. Extract 5-7 key insights from this transcript. "
        "Prioritise: surprising findings, counterintuitive ideas, and concrete actionable takeaways "
        "that would make someone stop scrolling. Avoid generic summaries. "
        "Each point must be 1 punchy sentence. "
        "Format as a plain list with one point per line (no numbers or bullets).\n\n"
        f"Transcript:\n{transcript[:4000]}...\n\n"
        "Key Points:"
    )
    resp = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[{"role": "user", "content": prompt}],
        max_tokens=400
    )
    points = [p.strip() for p in resp.choices[0].message.content.strip().split("\n") if p.strip()]
    return points[:7]


def extract_detailed_summary(transcript: str) -> str:
    """Generate 10-minute detailed summary."""
    client = _get_client()
    prompt = (
        "You are an expert at creating comprehensive yet concise summaries. "
        "Write a detailed 10-minute summary (~1500 words) of this video transcript. "
        "Organize with clear sections and subsections. Include specific examples and actionable insights.\n\n"
        f"Transcript:\n{transcript}\n\n"
        "Detailed Summary:"
    )
    resp = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[{"role": "user", "content": prompt}],
        max_tokens=1500
    )
    return resp.choices[0].message.content.strip()


def _extract_technical_chunk(chunk: str, idx: int) -> str:
    client = _get_client()
    prompt = (
        "Extract ALL technical content from this lecture segment with full fidelity. "
        "Your output will be used as the source for a spoken narration a student studies from. "
        "Include every: definition (formal + intuitive), theorem, equation, derivation step, "
        "algorithm, worked numerical example, proof sketch, condition, and edge case. "
        "Write equations in plain spoken English: e.g. 'the gradient of the loss with respect to w "
        "equals one over n times X transpose times the vector of residuals'. "
        "Reproduce every step of every derivation — never skip intermediate steps. "
        "Reproduce every numerical example in full. "
        "Skip only filler words, repeated questions, and off-topic tangents. "
        "Output as plain prose, one paragraph per topic, preserving the lecture's logical order.\n\n"
        f"Segment {idx + 1}:\n{chunk}"
    )
    resp = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[{"role": "user", "content": prompt}],
        max_tokens=2000,
    )
    return resp.choices[0].message.content.strip()


def _extract_technical_content(transcript: str, chunk_size: int = 2500) -> str:
    """Parallel extraction of verbatim technical content across the full transcript."""
    from concurrent.futures import ThreadPoolExecutor, as_completed
    words = transcript.split()
    chunks = [" ".join(words[i:i + chunk_size]) for i in range(0, len(words), chunk_size)]
    results = [None] * len(chunks)
    with ThreadPoolExecutor(max_workers=8) as executor:
        futures = {executor.submit(_extract_technical_chunk, chunk, i): i
                   for i, chunk in enumerate(chunks)}
        for future in as_completed(futures):
            results[futures[future]] = future.result()
    return "\n\n".join(r for r in results if r)


def generate_narration_script(transcript: str, key_points: List[str]) -> str:
    """Generate a technically complete, naturally spoken narration — like a brilliant TA reviewing the lecture."""
    import random
    client = _get_client()
    key_points_text = "\n".join(f"- {p}" for p in key_points)

    opening_styles = [
        "Start mid-story — drop the listener into a specific moment or scenario from the lecture before explaining anything.",
        "Open with the single most surprising or counterintuitive finding from the lecture. State it cold, no build-up.",
        "Open by naming the central problem or question the lecture is trying to solve, then hint at why the answer is non-obvious.",
        "Open with a vivid analogy that makes the core technical idea immediately intuitive.",
        "Open by stating what most students get wrong about this topic, then correct it.",
    ]
    chosen_opening = random.choice(opening_styles)

    # Extract verbatim technical content from the full transcript (equations, derivations, examples)
    print("Extracting technical content for narration brief...")
    technical_brief = _extract_technical_content(transcript)
    technical_ref = (
        "TECHNICAL REFERENCE (verbatim technical content extracted from the full lecture — "
        "every definition, equation, derivation step, and worked example):\n\n"
        + technical_brief
    )

    prompt = (
        "You are a brilliant, enthusiastic teaching assistant narrating a review of a lecture for students. "
        "Your voice is warm and clear, like a knowledgeable friend who genuinely loves the subject. "
        "Think of Richard Feynman explaining physics, or Andrej Karpathy walking through a concept on a whiteboard.\n\n"
        f"Opening instruction: {chosen_opening}\n\n"
        "Write a 15-20 minute narration script (~2500 words). Non-negotiable rules:\n\n"
        "TECHNICAL COMPLETENESS — this is the most important rule:\n"
        "- Every concept, formula, derivation, and example in the technical reference MUST appear in the narration.\n"
        "- When stating a formula or equation, read it out naturally: say 'the loss equals negative log of p of y given x' "
        "not just 'the loss function'. Spell out what each term means.\n"
        "- Do not skip steps. If a derivation has three steps, walk through all three in plain English.\n"
        "- Use the exact technical vocabulary from the reference. Do not replace jargon with vague language.\n"
        "- Reproduce any worked examples from the reference in the narration.\n\n"
        "SPOKEN STYLE — second priority:\n"
        "- Write entirely for the ear. No headers, bullets, numbered lists, markdown, or dashes as punctuation.\n"
        "- Never start with 'Have you ever', 'What if I told you', or 'In this video'.\n"
        "- Mix short punchy sentences with longer flowing ones for rhythm.\n"
        "- Use natural spoken connectors: 'Here is the thing.', 'Now notice what happens here.', "
        "'This is the key insight.', 'Think about what this means.'.\n"
        "- End with a synthesis: connect the ideas together and leave the student with a clear mental model.\n"
        "- Plain prose only. No special characters, no section labels.\n\n"
        f"Key concepts to cover (every single one must appear):\n{key_points_text}\n\n"
        f"{technical_ref}\n\n"
        "Narration script:"
    )
    resp = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[{"role": "user", "content": prompt}],
        max_tokens=5000
    )
    return resp.choices[0].message.content.strip()


def generate_clip_hooks(transcript: str, duration: float) -> List[Dict]:
    """Identify the 6 highest-signal moments in the transcript with precise timestamps."""
    client = _get_client()
    prompt = (
        "You are an expert video analyst and educator. Your job is to find the 6 most genuinely valuable moments "
        "in this transcript — not clickbait, but moments where the speaker says something that is:\n"
        "  - Counterintuitive or surprising relative to common belief\n"
        "  - A concrete example that makes an abstract idea click\n"
        "  - A specific number, study, or finding that changes how you think about the topic\n"
        "  - A turning point or reveal in the argument\n"
        "  - A practical takeaway someone could act on immediately\n\n"
        "Avoid: generic statements, introductions, transitions, or anything the audience already knows.\n\n"
        "For each moment output EXACTLY this format (no extra text):\n\n"
        "Title: [6-9 words — state the actual insight, not a teaser]\n"
        "Timestamp: [best estimate in seconds based on position in transcript]\n"
        "Description: [1-2 sentences explaining why this moment matters and what the viewer will learn]\n"
        "Overlay: [3-5 words that capture the core idea for on-screen text]\n"
        "Insight: [the single sentence from the transcript that best represents this moment]\n"
        "---\n\n"
        f"Video duration: {duration:.0f} seconds\n\n"
        f"Transcript:\n{transcript}\n\n"
        "Identify the 6 best moments:"
    )
    resp = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[{"role": "user", "content": prompt}],
        max_tokens=2000
    )

    text = resp.choices[0].message.content.strip()
    clips = []
    sections = text.split("---")
    for section in sections:
        section = section.strip()
        if not section:
            continue
        lines = section.split("\n")
        clip = {}
        desc_lines = []
        for line in lines:
            if line.startswith("Title:"):
                clip["title"] = line.replace("Title:", "").strip()
            elif line.startswith("Timestamp:"):
                try:
                    clip["timestamp"] = int(float(line.replace("Timestamp:", "").strip().split()[0]))
                except:
                    clip["timestamp"] = 0
            elif line.startswith("Description:"):
                desc_lines = [line.replace("Description:", "").strip()]
            elif line.startswith("Overlay:"):
                if desc_lines:
                    clip["description"] = " ".join(desc_lines)
                clip["overlay"] = line.replace("Overlay:", "").strip()
            elif line.startswith("Insight:"):
                clip["insight"] = line.replace("Insight:", "").strip()
            elif desc_lines and not any(line.startswith(k) for k in ("Title:", "Timestamp:", "Overlay:", "Insight:")):
                desc_lines.append(line.strip())
        if desc_lines and "description" not in clip:
            clip["description"] = " ".join(desc_lines)
        # Apply defaults for any missing optional fields so valid partial responses aren't dropped
        if "title" not in clip:
            continue
        clip.setdefault("timestamp", 0)
        clip.setdefault("description", clip.get("title", ""))
        clip.setdefault("overlay", " ".join(clip["title"].split()[:4]))
        clip.setdefault("insight", clip["description"])
        clip["timestamp"] = max(0, min(int(clip["timestamp"]), int(duration - 20)))
        clips.append(clip)

    # Fill any missing slots with evenly-spaced fallback clips
    clip_interval = max(1, int(duration / 6))
    for i in range(len(clips), 6):
        clips.append({
            "title": f"Key Moment {i+1}",
            "timestamp": min(i * clip_interval, int(duration - 20)),
            "description": f"Important segment #{i+1}",
            "overlay": f"Moment {i+1}",
            "insight": "Watch this key moment"
        })

    return clips[:6]


def generate_clips(video_path: str, working_dir: str, clip_hooks: List[Dict], clip_duration: int = 20) -> List[Dict]:
    """Generate MP4 clips from specified timestamps."""
    clips = []
    vid_dur = video_duration(video_path)
    for idx, hook in enumerate(clip_hooks):
        start = hook["timestamp"]
        start = max(0, min(start, vid_dur - clip_duration))
        # Sanitize title for use in filename — strip everything except alphanumerics and spaces
        safe_title = re.sub(r'[^a-zA-Z0-9 ]', '', hook['title']).strip().replace(' ', '_')[:20]
        outfile = os.path.join(working_dir, f"clip_{idx+1:02d}_{safe_title}.mp4")
        try:
            _run_cmd(["ffmpeg", "-y", "-ss", str(start), "-i", video_path, "-t", str(clip_duration), "-c:v", "libx264", "-c:a", "aac", "-preset", "ultrafast", outfile])
            clips.append({
                "clip_index": idx + 1,
                "title": hook["title"],
                "start": start,
                "duration": clip_duration,
                "path": outfile,
                "description": hook["description"],
                "overlay_text": hook["overlay"],
                "insight": hook["insight"]
            })
        except Exception as e:
            print(f"Warning: failed to generate clip {idx+1}: {e}")

    return clips


def generate_seo_metadata(transcript: str, key_points: List[str]) -> Dict:
    """Generate SEO-friendly metadata."""
    client = _get_client()
    prompt = (
        "You are an SEO expert for educational content. Based on this transcript, generate:\n"
        "Title: [catchy title, 5-8 words]\n"
        "Hashtags: [5 relevant hashtags, comma-separated]\n"
        "Description: [150-word YouTube description]\n\n"
        f"Transcript summary:\n{transcript[:1500]}...\n\n"
        "Key points:\n" + "\n".join(key_points) + "\n\n"
        "Output in the exact format above:"
    )
    resp = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[{"role": "user", "content": prompt}],
        max_tokens=500
    )
    
    text = resp.choices[0].message.content.strip()
    metadata = {}
    lines = text.split("\n")
    for line in lines:
        if line.startswith("Title:"):
            metadata["title"] = line.replace("Title:", "").strip()
        elif line.startswith("Hashtags:"):
            hashtags = [h.strip().lstrip("#") for h in line.replace("Hashtags:", "").strip().split(",")]
            metadata["hashtags"] = hashtags
        elif line.startswith("Description:"):
            metadata["description"] = line.replace("Description:", "").strip()
    
    # Fallback
    if not metadata.get("title"):
        metadata = {
            "title": "Educational Video Summary",
            "hashtags": ["education", "learning", "edtech"],
            "description": "Watch this comprehensive summary of the video."
        }
    
    return metadata


def _chunk_text(text: str, max_chars: int = 4000) -> List[str]:
    words = text.split()
    chunks = []
    current_chunk = []
    current_len = 0

    for word in words:
        if current_len + len(word) + 1 > max_chars:
            chunks.append(" ".join(current_chunk))
            current_chunk = [word]
            current_len = len(word) + 1
        else:
            current_chunk.append(word)
            current_len += len(word) + 1

    if current_chunk:
        chunks.append(" ".join(current_chunk))

    return chunks


def _filter_markdown_for_narration(text: str) -> str:
    # Code blocks
    text = re.sub(r'```[\s\S]*?```', '', text)
    text = re.sub(r'`[^`]+`', '', text)
    # Headings
    text = re.sub(r'^#{1,6}\s+', '', text, flags=re.MULTILINE)
    # Bold / italic / strikethrough — unwrap, keep inner text
    text = re.sub(r'(\*{1,3}|_{1,3})(.*?)\1', r'\2', text)
    text = re.sub(r'~~(.*?)~~', r'\1', text)
    # Horizontal rules (---, ***, ___, ===, any repeated symbol line)
    text = re.sub(r'^\s*([*\-_=])\s*(\1\s*){2,}$', '', text, flags=re.MULTILINE)
    # Links — keep display text
    text = re.sub(r'\[([^\]]+)\]\([^\)]*\)', r'\1', text)
    # Images — drop entirely
    text = re.sub(r'!\[[^\]]*\]\([^\)]*\)', '', text)
    # Blockquotes
    text = re.sub(r'^\s*>\s*', '', text, flags=re.MULTILINE)
    # List markers — convert to sentence flow
    text = re.sub(r'^\s*[-*+]\s+', '', text, flags=re.MULTILINE)
    text = re.sub(r'^\s*\d+[.)]\s+', '', text, flags=re.MULTILINE)
    # HTML tags
    text = re.sub(r'<[^>]+>', '', text)
    # Remaining non-speech symbols (pipes, backslashes, carets, tildes, etc.)
    text = re.sub(r'[|\\^~`@]', '', text)
    # Collapse whitespace
    text = re.sub(r'\n{3,}', '\n\n', text)
    text = re.sub(r' {2,}', ' ', text)
    return text.strip()


_VOICE_MAP = {
    "friendly":      "nova",
    "authoritative": "onyx",
    "energetic":     "shimmer",
}


def _ffmpeg_escape(text: str) -> str:
    text = text.replace("'", "\\'")
    text = text.replace(":", "\\:")
    text = text.replace("%", "\\%")
    return text


def generate_narration_audio(text: str, out_dir: str, voice_style: str = "friendly", filename: str = "summary_narration.mp3") -> str:
    clean_text = _filter_markdown_for_narration(text)
    voice = _VOICE_MAP.get((voice_style or "friendly").lower(), "nova")
    client = _get_client()

    # OpenAI TTS max 4096 chars per request — chunk by character count
    chunks = _chunk_text(clean_text, max_chars=4000)
    chunk_files = []

    for i, chunk in enumerate(chunks):
        chunk_path = os.path.join(out_dir, f"tts_chunk_{i:03d}.mp3")
        response = client.audio.speech.create(
            model="tts-1-hd",
            voice=voice,
            input=chunk,
            speed=1.15,
        )
        with open(chunk_path, "wb") as f:
            f.write(response.content)
        chunk_files.append(chunk_path)

    output_path = os.path.join(out_dir, filename)
    if len(chunk_files) == 1:
        os.rename(chunk_files[0], output_path)
        return output_path

    concat_list = os.path.join(out_dir, "tts_concat.txt")
    with open(concat_list, "w") as f:
        for p in chunk_files:
            f.write(f"file '{p}'\n")
    _run_cmd(["ffmpeg", "-y", "-f", "concat", "-safe", "0", "-i", concat_list,
              "-c", "copy", output_path])
    return output_path


def extract_keyframes(video_path: str, out_dir: str,
                       interval_secs: float = 15.0, max_frames: int = 60) -> List[str]:
    """Extract one frame every interval_secs using a single ffmpeg fps-filter call."""
    frames_dir = os.path.join(out_dir, "keyframes")
    os.makedirs(frames_dir, exist_ok=True)
    dur = video_duration(video_path)
    # Never exceed max_frames regardless of video length
    effective_interval = max(interval_secs, dur / max_frames)
    _run_cmd([
        "ffmpeg", "-y", "-i", video_path,
        "-vf", f"fps=1/{effective_interval:.2f},scale=1280:720",
        "-q:v", "4",
        os.path.join(frames_dir, "frame_%04d.jpg"),
    ])
    return sorted(str(p) for p in Path(frames_dir).glob("frame_*.jpg"))


def get_audio_duration(audio_path: str) -> float:
    out = _run_cmd(["ffprobe", "-v", "error", "-show_entries", "format=duration", "-of", "default=noprint_wrappers=1:nokey=1", audio_path])
    return float(out.strip())


def _render_slide(base_image_path: str, label: str, out_path: str, width: int, height: int):
    """Composite a text label onto a keyframe using Pillow and save as JPEG."""
    if base_image_path and os.path.exists(base_image_path):
        img = Image.open(base_image_path).convert("RGB")
        img = img.resize((width, height), Image.LANCZOS)
    else:
        img = Image.new("RGB", (width, height), (0, 0, 0))

    draw = ImageDraw.Draw(img)

    # Try to load a system font, fall back to default
    font_large, font_small = None, None
    for font_path in [
        "/System/Library/Fonts/Helvetica.ttc",
        "/System/Library/Fonts/Arial.ttf",
        "/Library/Fonts/Arial.ttf",
    ]:
        if os.path.exists(font_path):
            try:
                font_large = ImageFont.truetype(font_path, 52)
                font_small = ImageFont.truetype(font_path, 34)
            except Exception:
                pass
            break
    if font_large is None:
        font_large = ImageFont.load_default()
        font_small = font_large

    # Dark gradient bar at the bottom for readability
    overlay = Image.new("RGBA", (width, height), (0, 0, 0, 0))
    bar = ImageDraw.Draw(overlay)
    bar.rectangle([(0, height - 180), (width, height)], fill=(0, 0, 0, 160))
    img = Image.alpha_composite(img.convert("RGBA"), overlay).convert("RGB")
    draw = ImageDraw.Draw(img)

    # Wrap and draw text inside the bar
    wrapped = textwrap.fill(label, width=52)
    lines = wrapped.split("\n")
    font = font_large if len(lines) == 1 else font_small
    font_size = font.size if hasattr(font, "size") else 34
    line_height = font_size + 8
    total_h = len(lines) * line_height
    y = height - 160 + (160 - total_h) // 2
    for line in lines:
        bbox = draw.textbbox((0, 0), line, font=font)
        x = (width - (bbox[2] - bbox[0])) // 2
        draw.text((x, y), line, font=font, fill=(255, 255, 255))
        y += line_height

    img.save(out_path, "JPEG", quality=90)



def _render_text_overlay(label: str, out_path: str, width: int = 1280, height: int = 720):
    """Render a transparent RGBA PNG with gradient bar + subtitle-style text for ffmpeg overlay."""
    overlay = Image.new("RGBA", (width, height), (0, 0, 0, 0))
    draw = ImageDraw.Draw(overlay)

    # Load font at subtitle size
    font = None
    for font_path in [
        "/System/Library/Fonts/Helvetica.ttc",
        "/System/Library/Fonts/Arial.ttf",
        "/Library/Fonts/Arial.ttf",
    ]:
        if os.path.exists(font_path):
            try:
                font = ImageFont.truetype(font_path, 30)
            except Exception:
                pass
            break
    if font is None:
        font = ImageFont.load_default()

    # Pixel-accurate word wrap — never overflows the frame
    h_margin = 80  # px from each side
    max_text_width = width - 2 * h_margin
    words = label.split()
    lines, current = [], []
    for word in words:
        test = " ".join(current + [word])
        w = draw.textbbox((0, 0), test, font=font)[2]
        if w <= max_text_width:
            current.append(word)
        else:
            if current:
                lines.append(" ".join(current))
            current = [word]
    if current:
        lines.append(" ".join(current))

    # Cap at 3 lines — captions should be glanceable
    lines = lines[:3]

    font_size = font.size if hasattr(font, "size") else 30
    line_height = font_size + 8
    v_padding = 18
    bar_h = len(lines) * line_height + v_padding * 2
    bar_top = height - bar_h - 20  # 20px from bottom edge

    # Gradient bar — transparent at top, semi-opaque at bottom
    for row in range(bar_h):
        alpha = int(200 * (row / bar_h))
        draw.line([(0, bar_top + row), (width, bar_top + row)], fill=(0, 0, 0, alpha))

    # Draw each line centred
    y = bar_top + v_padding
    for line in lines:
        bbox = draw.textbbox((0, 0), line, font=font)
        x = (width - (bbox[2] - bbox[0])) // 2
        draw.text((x + 1, y + 1), line, font=font, fill=(0, 0, 0, 200))   # shadow
        draw.text((x, y), line, font=font, fill=(255, 255, 255, 255))
        y += line_height

    overlay.save(out_path, "PNG")


def create_summary_video(narration_audio_path: str, out_dir: str, key_points: List[str] = None,
                          frame_paths: List[str] = None, narration_text: str = None,
                          width: int = 1280, height: int = 720) -> str:
    duration = get_audio_duration(narration_audio_path)
    video_path = os.path.join(out_dir, "summary_video.mp4")

    if not key_points:
        key_points = []

    # Strip markdown so it never appears on screen
    clean_narration = _filter_markdown_for_narration(narration_text) if narration_text else None

    # Build segment list: (label, keyframe_path_or_None)
    frames = frame_paths or []
    if clean_narration:
        caption_segs = _split_into_caption_segments(clean_narration, duration, target_secs=12.0)
        slides = [
            (seg["text"], frames[i % len(frames)] if frames else None)
            for i, seg in enumerate(caption_segs)
        ]
        seg_durations = [seg["end"] - seg["start"] for seg in caption_segs]
    else:
        slides = [("Video Summary", frames[0] if frames else None)]
        for idx, pt in enumerate((key_points or [])[:6]):
            slides.append((pt, frames[idx % len(frames)] if frames else None))
        slides.append(("Thanks for watching!", frames[-1] if frames else None))
        seg_dur = duration / len(slides)
        seg_durations = [seg_dur] * len(slides)

    # Still keyframe + text overlay, fade in at start and fade out at end for smooth transitions
    FADE = 0.4  # seconds for fade in/out
    seg_paths = []
    for i, (label, frame) in enumerate(slides):
        overlay_path = os.path.join(out_dir, f"overlay_{i:03d}.png")
        _render_text_overlay(label, overlay_path, width, height)

        if frame and os.path.exists(frame):
            bg_path = frame
        else:
            fallback = os.path.join(out_dir, f"bg_fallback_{i:03d}.jpg")
            Image.new("RGB", (width, height), (15, 15, 30)).save(fallback, "JPEG")
            bg_path = fallback

        dur = seg_durations[i]
        fade_d = min(FADE, dur / 3)  # never let fade eat more than 1/3 of a short segment
        fade_out_start = max(0, dur - fade_d)

        # Scale bg, overlay text PNG, then fade in + fade out
        filter_complex = (
            f"[0:v]scale={width}:{height}[bg];"
            f"[bg][1:v]overlay=0:0,"
            f"fade=t=in:st=0:d={fade_d:.3f},"
            f"fade=t=out:st={fade_out_start:.3f}:d={fade_d:.3f}[out]"
        )
        seg_path = os.path.join(out_dir, f"seg_{i:03d}.mp4")
        _run_cmd([
            "ffmpeg", "-y",
            "-loop", "1", "-i", bg_path,
            "-loop", "1", "-i", overlay_path,
            "-filter_complex", filter_complex,
            "-map", "[out]",
            "-t", f"{dur:.3f}",
            "-c:v", "libx264", "-r", "25", "-pix_fmt", "yuv420p",
            seg_path
        ])
        seg_paths.append(seg_path)

    # Step 3: Concatenate all segments (concat demuxer — one call, no filter_complex limit)
    concat_list = os.path.join(out_dir, "seg_concat.txt")
    with open(concat_list, "w") as f:
        for p in seg_paths:
            f.write(f"file '{p}'\n")
    intermediate_path = os.path.join(out_dir, "intermediate.mp4")
    _run_cmd([
        "ffmpeg", "-y", "-f", "concat", "-safe", "0", "-i", concat_list,
        "-c", "copy", intermediate_path
    ])

    # Step 4: Mix in narration audio
    _run_cmd([
        "ffmpeg", "-y",
        "-i", intermediate_path,
        "-i", narration_audio_path,
        "-c:v", "copy", "-c:a", "aac", "-b:a", "192k",
        "-shortest", video_path
    ])

    return video_path


def _split_into_caption_segments(text: str, audio_duration: float, target_secs: float = 12.0) -> List[Dict]:
    """Split narration text into timed segments proportional to word count."""
    words = text.split()
    if not words:
        return [{"text": "", "start": 0.0, "end": audio_duration}]
    words_per_sec = len(words) / audio_duration
    target_words = max(1, int(target_secs * words_per_sec))
    segments = []
    for i in range(0, len(words), target_words):
        chunk = words[i:i + target_words]
        start = i / words_per_sec
        end = min((i + len(chunk)) / words_per_sec, audio_duration)
        segments.append({"text": " ".join(chunk), "start": start, "end": end})
    return segments


def generate_highlight_reel(clips: List[Dict], out_dir: str, max_clips: int = 3) -> str:
    """Concatenate the top clips into a 30s-1min highlight reel."""
    reel_path = os.path.join(out_dir, "highlight_reel.mp4")
    selected = [c for c in clips if os.path.exists(c["path"])][:max_clips]
    if not selected:
        return ""
    concat_list = os.path.join(out_dir, "reel_concat.txt")
    with open(concat_list, "w") as f:
        for clip in selected:
            f.write(f"file '{clip['path']}'\n")
    _run_cmd(["ffmpeg", "-y", "-f", "concat", "-safe", "0", "-i", concat_list, "-c", "copy", reel_path])
    return reel_path


def _slides_from_chunk(chunk: str, chunk_index: int, total_chunks: int) -> List[Dict]:
    """Generate slides from one transcript chunk — one GPT call per chunk."""
    client = _get_client()
    prompt = (
        "You are a PhD-level teaching assistant creating a comprehensive exam reference for university students. "
        "Your output will be the ONLY resource a student has to prepare for their exam on this lecture. "
        "It must be complete enough that a student who never attended could fully master the material.\n\n"
        f"This is chunk {chunk_index + 1} of {total_chunks} from the lecture transcript.\n\n"
        "Produce one slide for EACH distinct concept, definition, theorem, algorithm, derivation step, "
        "worked example, proof, edge case, or important remark in this chunk. "
        "When in doubt, make more slides — never merge two ideas onto one slide.\n\n"
        "Each slide is a JSON object with exactly these fields:\n\n"
        "  title: precise heading (5-10 words) naming the specific concept — not generic labels like 'Introduction'\n\n"
        "  concept: EXHAUSTIVE explanation — write as much as needed. Include:\n"
        "    - Formal definition using exact notation and terminology from the lecture\n"
        "    - Intuitive explanation of WHY it works or what it means geometrically/physically\n"
        "    - All assumptions, conditions, and constraints that must hold\n"
        "    - How it connects to preceding concepts in the lecture\n"
        "    - Common misconceptions or edge cases the lecturer highlighted\n"
        "    - Any motivating example or real-world context given\n"
        "    Minimum 5-8 sentences. Do not truncate. Self-contained — no 'as mentioned earlier'.\n\n"
        "  detail: ALL formulas, equations, derivations, algorithm pseudocode, and worked numerical examples "
        "from the lecture for this concept. Write every step — do not skip intermediate steps in derivations. "
        "Use plain text math notation: e.g. 'grad_w L = (1/n) * X^T * (sigma(Xw) - y)'. "
        "If the lecturer worked through a numerical example, reproduce every calculation. "
        "If there are multiple formulas, include all of them separated by newlines. "
        "Empty string only if the concept is genuinely non-mathematical.\n\n"
        "  takeaway: one precise exam-ready sentence — the single most important thing to remember, "
        "including the key condition or caveat that distinguishes this from similar concepts.\n\n"
        "Absolute rules:\n"
        "- Use only content present in this chunk. Do not invent or hallucinate.\n"
        "- Use the lecturer's exact variable names, notation, and terminology.\n"
        "- Never truncate a derivation or example partway through.\n"
        "- A worked example is its own slide, not merged into the concept slide.\n"
        "Return a JSON array only. No markdown, no other text.\n\n"
        f"Transcript chunk:\n{chunk}"
    )
    resp = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[{"role": "user", "content": prompt}],
        max_tokens=4000,
    )
    raw = resp.choices[0].message.content.strip()
    match = re.search(r'\[[\s\S]*\]', raw)
    if match:
        try:
            return json.loads(match.group(0))
        except Exception:
            pass
    return []


def _generate_slide_content(transcript: str, exec_summary: str, key_points: List[str]) -> List[Dict]:
    """Chunk the full transcript and generate slides per chunk — parallel GPT calls."""
    from concurrent.futures import ThreadPoolExecutor, as_completed
    words = transcript.split()
    chunk_size = 2500
    chunks = [" ".join(words[i:i + chunk_size]) for i in range(0, len(words), chunk_size)]
    chunks = chunks[:8]

    results = [None] * len(chunks)
    with ThreadPoolExecutor(max_workers=8) as executor:
        futures = {executor.submit(_slides_from_chunk, chunk, i, len(chunks)): i
                   for i, chunk in enumerate(chunks)}
        for future in as_completed(futures):
            results[futures[future]] = future.result()

    all_slides = [slide for r in results if r for slide in r]
    if not all_slides:
        all_slides = [{"title": p[:80], "concept": p, "detail": "", "takeaway": ""} for p in key_points]
    return all_slides


def generate_slide_presentation(
    out_dir: str,
    key_points: List[str],
    exec_summary: str,
    frame_paths: List[str],
    clip_hooks: List[Dict],
    seo: Dict,
    slide_data: List[Dict] = None,
    transcript: str = "",
) -> str:
    """Generate a clean, content-rich .pptx slide deck."""

    # Palette — clean dark theme
    C_BG        = RGBColor(18,  18,  32)   # near-black navy
    C_ACCENT    = RGBColor(99,  179, 237)  # sky blue
    C_WHITE     = RGBColor(245, 245, 250)
    C_SUBTEXT   = RGBColor(160, 170, 190)
    C_DIVIDER   = RGBColor(50,  55,  80)
    C_DETAIL_BG = RGBColor(28,  32,  52)   # slightly lighter panel

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    W = prs.slide_width
    H = prs.slide_height

    # ── helpers ────────────────────────────────────────────────────────────

    def _bg(slide):
        """Solid dark background."""
        s = slide.shapes.add_shape(1, 0, 0, W, H)
        s.fill.solid()
        s.fill.fore_color.rgb = C_BG
        s.line.fill.background()

    def _rect(slide, left, top, width, height, color):
        s = slide.shapes.add_shape(1, left, top, width, height)
        s.fill.solid()
        s.fill.fore_color.rgb = color
        s.line.fill.background()
        return s

    def _txb(slide, text, left, top, width, height,
             size=20, bold=False, color=None, align=PP_ALIGN.LEFT, wrap=True):
        color = color or C_WHITE
        txb = slide.shapes.add_textbox(left, top, width, height)
        tf  = txb.text_frame
        tf.word_wrap = wrap
        p   = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size  = Pt(size)
        run.font.bold  = bold
        run.font.color.rgb = color
        return txb

    def _bullets(slide, items, left, top, width, height, size=16):
        """Render a list of strings as bullet paragraphs inside one textbox."""
        txb = slide.shapes.add_textbox(left, top, width, height)
        tf  = txb.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = f"•  {item}"
            run.font.size  = Pt(size)
            run.font.color.rgb = C_WHITE

    def _accent_bar(slide, top, height=Inches(0.06)):
        """Thin horizontal accent line."""
        _rect(slide, Inches(0.5), top, W - Inches(1.0), height, C_ACCENT)

    # ── Slide 1: Title + Executive Summary ─────────────────────────────────
    slide = prs.slides.add_slide(blank)
    _bg(slide)
    # left accent stripe
    _rect(slide, 0, 0, Inches(0.18), H, C_ACCENT)
    # title
    _txb(slide, seo.get("title", "Video Summary"),
         Inches(0.5), Inches(1.6), W - Inches(1.0), Inches(1.4),
         size=38, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
    _accent_bar(slide, Inches(3.1))
    # exec summary body
    _txb(slide, exec_summary,
         Inches(0.5), Inches(3.3), W - Inches(1.0), Inches(3.2),
         size=19, color=C_SUBTEXT, align=PP_ALIGN.LEFT)
    # hashtags bottom-right
    tags = "  ".join(f"#{t}" for t in seo.get("hashtags", [])[:5])
    _txb(slide, tags, Inches(0.5), H - Inches(0.55), W - Inches(1.0), Inches(0.45),
         size=12, color=C_ACCENT, align=PP_ALIGN.LEFT)

    # Use pre-generated slide_data if provided, otherwise generate now
    if not slide_data:
        slide_data = _generate_slide_content(transcript, exec_summary, key_points)

    # ── Slides 2-N: One concept per slide ──────────────────────────────────
    total_content_slides = len(slide_data)
    for idx, sd in enumerate(slide_data):
        slide = prs.slides.add_slide(blank)
        _bg(slide)

        # slide number tag top-right (slide 2 onwards, 1 = title)
        _txb(slide, f"{idx+2} / {total_content_slides + 1}",
             W - Inches(1.4), Inches(0.2), Inches(1.2), Inches(0.4),
             size=11, color=C_SUBTEXT, align=PP_ALIGN.RIGHT)

        # accent bar top
        _rect(slide, 0, 0, W, Inches(0.12), C_ACCENT)

        # slide title
        _txb(slide, sd.get("title", f"Concept {idx+1}"),
             Inches(0.5), Inches(0.3), W - Inches(1.5), Inches(0.8),
             size=28, bold=True, color=C_WHITE)

        # thin divider
        _accent_bar(slide, Inches(1.2), height=Inches(0.03))

        detail_text = sd.get("detail", "").strip()

        if detail_text:
            # Two-column: concept left, detail panel right
            col_w = W / 2 - Inches(0.7)
            _txb(slide, "EXPLANATION",
                 Inches(0.5), Inches(1.35), col_w, Inches(0.35),
                 size=11, bold=True, color=C_ACCENT)
            _txb(slide, sd.get("concept", ""),
                 Inches(0.5), Inches(1.75), col_w, Inches(3.8),
                 size=17, color=C_WHITE)

            rx = W / 2 + Inches(0.1)
            _rect(slide, rx - Inches(0.15), Inches(1.35),
                  col_w + Inches(0.3), Inches(4.3), C_DETAIL_BG)
            _txb(slide, "FORMULA / EXAMPLE",
                 rx, Inches(1.45), col_w, Inches(0.35),
                 size=11, bold=True, color=C_ACCENT)
            _txb(slide, detail_text,
                 rx, Inches(1.85), col_w, Inches(3.6),
                 size=16, color=C_WHITE)
        else:
            # Full width — concept gets all the space
            _txb(slide, sd.get("concept", ""),
                 Inches(0.5), Inches(1.35), W - Inches(1.0), Inches(4.8),
                 size=19, color=C_WHITE)

        # takeaway strip at bottom
        takeaway = sd.get("takeaway", "")
        if takeaway:
            _rect(slide, 0, H - Inches(0.85), W, Inches(0.85), C_DETAIL_BG)
            _txb(slide, "KEY TAKEAWAY",
                 Inches(0.5), H - Inches(0.82), Inches(2.2), Inches(0.35),
                 size=11, bold=True, color=C_ACCENT)
            _txb(slide, takeaway,
                 Inches(2.5), H - Inches(0.82), W - Inches(3.0), Inches(0.75),
                 size=15, color=C_WHITE)

    # ── Final slide: Clip moments ──────────────────────────────────────────
    if clip_hooks:
        slide = prs.slides.add_slide(blank)
        _bg(slide)
        _rect(slide, 0, 0, W, Inches(0.12), C_ACCENT)
        _txb(slide, "Moments Worth Revisiting",
             Inches(0.5), Inches(0.3), W - Inches(1.0), Inches(0.8),
             size=28, bold=True, color=C_WHITE)
        _accent_bar(slide, Inches(1.2), height=Inches(0.03))

        items = [
            f"{h.get('title','')} — {h.get('insight','')}"
            for h in clip_hooks[:6]
        ]
        _bullets(slide, items, Inches(0.5), Inches(1.4),
                 W - Inches(1.0), Inches(5.5), size=17)

    pptx_path = os.path.join(out_dir, "summary_slides.pptx")
    prs.save(pptx_path)
    return pptx_path


def generate_storyboard(out_dir: str, duration: float, key_points: List[str], clip_hooks: List[Dict], voice_style: str = "friendly") -> str:
    storyboard = []
    # intro slide
    storyboard.append({
        "start": 0,
        "end": 3,
        "type": "intro",
        "text": "Video ShortSage 10-minute summary: key ideas and clips",
        "image_placeholder": "intro_slide.png"
    })

    # key points slides (5 seconds each)
    for i, kp in enumerate(key_points[:6]):
        start = 3 + i * 5
        end = min(start + 5, duration)
        storyboard.append({
            "start": start,
            "end": end,
            "type": "key_point",
            "text": kp,
            "image_placeholder": f"keypoint_{i+1}.png"
        })

    # clips section metadata
    for i, clip in enumerate(clip_hooks[:6]):
        storyboard.append({
            "start": clip.get("timestamp", 0),
            "end": min(clip.get("timestamp", 0) + 20, duration),
            "type": "clip_hook",
            "title": clip.get("title", "Clip"),
            "description": clip.get("description", ""),
            "overlay": clip.get("overlay", ""),
            "image_placeholder": f"clip_{i+1}.png"
        })

    # outro slide
    final_start = max(duration - 3, 0)
    storyboard.append({
        "start": final_start,
        "end": duration,
        "type": "outro",
        "text": "Thank you for watching! Follow for more educational summaries.",
        "image_placeholder": "outro_slide.png"
    })

    storyboard_path = os.path.join(out_dir, "storyboard.json")
    with open(storyboard_path, "w") as f:
        json.dump({
            "duration": duration,
            "voice_style": voice_style,
            "storyboard": storyboard
        }, f, indent=2)

    return storyboard_path


def generate_pdf_slides(slide_data: List[Dict], out_dir: str, seo: Dict, exec_summary: str) -> str:
    """Render slide_data as a multi-page PDF using Pillow — no extra dependencies."""
    PW, PH = 1280, 720
    BG      = (18,  18,  32)
    ACCENT  = (99,  179, 237)
    WHITE   = (245, 245, 250)
    SUBTEXT = (160, 170, 190)
    PANEL   = (28,  32,  52)
    DIVIDER = (50,  55,  80)

    # Load fonts once
    font_h = font_body = font_small = font_tiny = None
    for fp in ["/System/Library/Fonts/Helvetica.ttc", "/System/Library/Fonts/Arial.ttf", "/Library/Fonts/Arial.ttf"]:
        if os.path.exists(fp):
            try:
                font_h     = ImageFont.truetype(fp, 36)
                font_body  = ImageFont.truetype(fp, 19)
                font_small = ImageFont.truetype(fp, 16)
                font_tiny  = ImageFont.truetype(fp, 12)
            except Exception:
                pass
            break
    if font_h is None:
        font_h = font_body = font_small = font_tiny = ImageFont.load_default()

    def _wrap(draw, text, font, max_w):
        words = text.split()
        lines, cur = [], []
        for w in words:
            test = " ".join(cur + [w])
            if draw.textbbox((0, 0), test, font=font)[2] <= max_w:
                cur.append(w)
            else:
                if cur:
                    lines.append(" ".join(cur))
                cur = [w]
        if cur:
            lines.append(" ".join(cur))
        return lines

    def _draw_text_block(draw, text, font, x, y, max_w, color, max_y=None):
        lh = (font.size if hasattr(font, "size") else 17) + 6
        for line in _wrap(draw, text, font, max_w):
            if max_y and y + lh > max_y:
                break
            draw.text((x, y), line, font=font, fill=color)
            y += lh
        return y

    images = []

    # ── Title slide ──────────────────────────────────────────────────────────
    img  = Image.new("RGB", (PW, PH), BG)
    draw = ImageDraw.Draw(img)
    draw.rectangle([(0, 0), (14, PH)], fill=ACCENT)          # left stripe
    _draw_text_block(draw, seo.get("title", "Lecture Summary"), font_h, 36, 100, PW - 72, WHITE)
    draw.rectangle([(36, 200), (PW - 36, 203)], fill=ACCENT)  # divider
    _draw_text_block(draw, exec_summary, font_small, 36, 220, PW - 72, SUBTEXT, max_y=PH - 50)
    images.append(img)

    # ── Content slides ────────────────────────────────────────────────────────
    for idx, sd in enumerate(slide_data):
        img  = Image.new("RGB", (PW, PH), BG)
        draw = ImageDraw.Draw(img)
        draw.rectangle([(0, 0), (PW, 8)], fill=ACCENT)         # top bar

        # slide number
        num_text = f"{idx + 2} / {len(slide_data) + 1}"
        nw = draw.textbbox((0, 0), num_text, font=font_tiny)[2]
        draw.text((PW - nw - 20, 16), num_text, font=font_tiny, fill=SUBTEXT)

        # title
        _draw_text_block(draw, sd.get("title", ""), font_h, 30, 20, PW - 160, WHITE)
        draw.rectangle([(30, 80), (PW - 30, 82)], fill=DIVIDER)

        detail  = sd.get("detail",   "").strip()
        concept = sd.get("concept",  "")
        takeaway = sd.get("takeaway", "")
        bottom  = PH - 70 if takeaway else PH - 20

        if detail:
            col_w = PW // 2 - 50
            draw.text((30, 92), "EXPLANATION", font=font_tiny, fill=ACCENT)
            _draw_text_block(draw, concept, font_small, 30, 112, col_w, WHITE, max_y=bottom)
            rx = PW // 2 + 10
            draw.rectangle([(rx - 8, 88), (PW - 20, bottom + 4)], fill=PANEL)
            draw.text((rx, 92), "FORMULA / EXAMPLE", font=font_tiny, fill=ACCENT)
            _draw_text_block(draw, detail, font_small, rx, 112, col_w, WHITE, max_y=bottom)
        else:
            _draw_text_block(draw, concept, font_body, 30, 96, PW - 60, WHITE, max_y=bottom)

        if takeaway:
            draw.rectangle([(0, PH - 66), (PW, PH)], fill=PANEL)
            draw.text((30, PH - 60), "KEY TAKEAWAY", font=font_tiny, fill=ACCENT)
            _draw_text_block(draw, takeaway, font_small, 210, PH - 60, PW - 240, WHITE)

        images.append(img)

    pdf_path = os.path.join(out_dir, "summary_slides.pdf")
    images[0].save(pdf_path, "PDF", save_all=True, append_images=images[1:])
    return pdf_path


def save_transcript(transcript: str, out_dir: str) -> str:
    """Save the full transcript as a plain .txt file."""
    path = os.path.join(out_dir, "transcript.txt")
    with open(path, "w", encoding="utf-8") as f:
        f.write(transcript)
    return path


def process(video_path: str, working_dir: str, voice_style: str = "friendly") -> Dict:
    """Main processing pipeline: transcript-only, output-optimized."""
    print(f"Processing {video_path} with voice style {voice_style}...")
    
    # Extract audio & transcribe
    print("Extracting audio...")
    audio_path = extract_audio(video_path, working_dir)
    
    print("Transcribing speech...")
    transcript = speech_to_text(audio_path)
    
    duration = video_duration(video_path)

    # Stage 1: Run all independent analysis calls + keyframe extraction in parallel
    print("Running parallel analysis (summary, key points, clip hooks, keyframes)...")
    from concurrent.futures import ThreadPoolExecutor
    with ThreadPoolExecutor(max_workers=5) as executor:
        f_exec    = executor.submit(extract_executive_summary, transcript)
        f_kp      = executor.submit(extract_key_points, transcript)
        f_hooks   = executor.submit(generate_clip_hooks, transcript, duration)
        f_frames  = executor.submit(extract_keyframes, video_path, working_dir, 15.0, 60)
        exec_summary = f_exec.result()
        key_points   = f_kp.result()
        clip_hooks   = f_hooks.result()
        frame_paths  = f_frames.result()

    # Stage 2: clip generation + SEO + slide content all run in parallel
    print("Running parallel generation (clips, SEO, slides)...")
    with ThreadPoolExecutor(max_workers=3) as executor:
        f_clips = executor.submit(generate_clips, video_path, working_dir, clip_hooks, 20)
        f_seo   = executor.submit(generate_seo_metadata, transcript, key_points)
        f_slides = executor.submit(_generate_slide_content, transcript, exec_summary, key_points)
        clips      = f_clips.result()
        seo        = f_seo.result()
        slide_data = f_slides.result()

    # Stage 3: narration (depends on slide_data)
    print("Generating narration script...")
    narration_script = generate_narration_script(transcript, key_points)

    print("Generating narration audio...")
    narration_audio = generate_narration_audio(narration_script, working_dir, voice_style=voice_style)

    # Stage 4: summary video + slides/pdf/transcript all run in parallel
    print("Running parallel finalization (video, slides, PDF, transcript)...")
    with ThreadPoolExecutor(max_workers=4) as executor:
        f_video = executor.submit(create_summary_video, narration_audio, working_dir,
                                  key_points, frame_paths, narration_script)
        f_pptx  = executor.submit(generate_slide_presentation, working_dir, key_points,
                                  exec_summary, frame_paths, clip_hooks, seo, slide_data)
        f_pdf   = executor.submit(generate_pdf_slides, slide_data, working_dir, seo, exec_summary)
        f_txt   = executor.submit(save_transcript, transcript, working_dir)
        summary_video_path   = f_video.result()
        slides_path          = f_pptx.result()
        pdf_path             = f_pdf.result()
        transcript_txt_path  = f_txt.result()

    print("Generating storyboard JSON...")
    storyboard_path = generate_storyboard(working_dir, duration, key_points, clip_hooks, voice_style=voice_style)

    return {
        "summary_video_path": summary_video_path,
        "summary_narration_audio": narration_audio,
        "slides_path": slides_path,
        "pdf_path": pdf_path,
        "transcript_txt_path": transcript_txt_path,
        "storyboard_path": storyboard_path,
        "source_video": video_path,
        "duration_seconds": duration,
        "transcript": transcript,
        "summaries": {
            "executive": exec_summary,
            "detailed": narration_script,
            "key_points": key_points
        },
        "clip_hooks": clip_hooks,
        "generated_clips": clips,
        "seo_metadata": seo,
        "publication_ready": {
            "short_form_content": [
                {
                    "platform": "TikTok/YouTube Shorts",
                    "clip_id": clip["clip_index"],
                    "title": clip["title"],
                    "description": clip["description"],
                    "overlay": clip["overlay_text"],
                    "file": os.path.basename(clip["path"])
                }
                for clip in clips
            ],
            "hashtags": seo.get("hashtags", []),
            "seo_title": seo.get("title", ""),
            "seo_description": seo.get("description", "")
        }
    }
